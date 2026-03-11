"""
PPTX 產生模組
支援橫式 / 直式投影片自動偵測，
若同一批投影片中包含混合方向，依方案 A 分割輸出為多個 .pptx 檔案。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from lib.config import get_slide_dimensions, is_portrait


def _build_presentation(slides_subset, slide_w_inches, slide_h_inches):
    """
    建立單一 Presentation 物件，包含 slides_subset 中所有頁面。
    slide_w_inches / slide_h_inches 為本份簡報統一使用的尺寸。
    """
    prs = Presentation()
    prs.slide_width = Inches(slide_w_inches)
    prs.slide_height = Inches(slide_h_inches)
    blank_layout = prs.slide_layouts[6]

    slide_h_pt = slide_h_inches * 72  # 換算為 points，供字型大小計算使用

    for slide_data in slides_subset:
        slide = prs.slides.add_slide(blank_layout)

        # 1. 背景圖片
        if slide_data.get('background_image'):
            try:
                slide.shapes.add_picture(
                    slide_data['background_image'],
                    0, 0,
                    prs.slide_width, prs.slide_height
                )
            except Exception as e:
                print(f"  [pptx_generator] 背景圖片加入失敗: {e}")

        # 2. 文字區塊
        for block in slide_data.get('text_blocks', []):
            text = block.get('text', '')
            if not text:
                continue

            # box_2d: [ymin, xmin, ymax, xmax]，座標比例 0-1000
            ymin, xmin, ymax, xmax = block.get('box_2d', [0, 0, 0, 0])

            left   = (xmin / 1000) * prs.slide_width
            top    = (ymin / 1000) * prs.slide_height
            width  = ((xmax - xmin) / 1000) * prs.slide_width
            height = ((ymax - ymin) / 1000) * prs.slide_height

            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = text

            # 字型大小：依實際投影片高度動態計算（非硬編碼）
            box_height_rel = (ymax - ymin) / 1000
            calculated_size = int(box_height_rel * slide_h_pt * 0.75)
            calculated_size = max(8, min(100, calculated_size))
            p.font.size = Pt(calculated_size)

            if block.get('is_bold'):
                p.font.bold = True

            color_hex = block.get('color', '000000').replace('#', '')
            try:
                p.font.color.rgb = RGBColor.from_string(color_hex)
            except Exception:
                p.font.color.rgb = RGBColor(0, 0, 0)

            align = block.get('align', 'left').lower()
            if 'center' in align:
                p.alignment = PP_ALIGN.CENTER
            elif 'right' in align:
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT

    return prs


def create_pptx(slides_data, output_file):
    """
    將投影片資料輸出為 PPTX 檔案。

    若 slides_data 中包含混合方向（橫式 + 直式），
    依方案 A 分割為多個 .pptx 檔案：
      - {base}_landscape.pptx（橫式頁面）
      - {base}_portrait.pptx（直式頁面）
    若只有單一方向，直接輸出為 output_file。

    每個 slide_data dict 可包含：
      - background_image: BytesIO 或檔案路徑（選填）
      - text_blocks: OCR 區塊列表（選填）
      - page_size: (width_pt, height_pt)（選填，未提供時預設為橫式）

    回傳：實際儲存的檔案路徑列表。
    """
    # 為每頁決定方向
    orientations = []
    for sd in slides_data:
        pw, ph = sd.get('page_size', (1, 0))  # 預設寬 > 高 = 橫式
        orientations.append('portrait' if is_portrait(pw, ph) else 'landscape')

    unique_orientations = set(orientations)

    if len(unique_orientations) == 1:
        # 單一方向：直接輸出
        orientation = unique_orientations.pop()
        pw, ph = slides_data[0].get('page_size', (1, 0))
        slide_w, slide_h = get_slide_dimensions(pw, ph)
        prs = _build_presentation(slides_data, slide_w, slide_h)
        prs.save(output_file)
        print(f"PPTX 已儲存: {output_file}")
        return [output_file]

    else:
        # 混合方向：分割輸出（方案 A）
        print("  [pptx_generator] 偵測到混合方向頁面，將分割輸出多個 PPTX 檔案")
        base, ext = os.path.splitext(output_file)
        saved_files = []

        for orientation in ('landscape', 'portrait'):
            subset = [
                sd for sd, o in zip(slides_data, orientations)
                if o == orientation
            ]
            if not subset:
                continue

            pw, ph = subset[0].get('page_size', (1, 0))
            slide_w, slide_h = get_slide_dimensions(pw, ph)
            prs = _build_presentation(subset, slide_w, slide_h)

            label = '橫式' if orientation == 'landscape' else '直式'
            out_path = f"{base}_{orientation}{ext}"
            prs.save(out_path)
            print(f"PPTX 已儲存 ({label}，共 {len(subset)} 頁): {out_path}")
            saved_files.append(out_path)

        return saved_files
